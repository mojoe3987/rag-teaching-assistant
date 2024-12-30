from pathlib import Path
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_from_pdf_per_page(file_path):
    """
    Returns a list of (page_text, metadata) for each page in the PDF.
    """
    reader = PdfReader(file_path)
    docs = []
    # Optionally parse session number from filename (e.g. 'session_7.pdf')
    session_num = parse_session_from_filename(file_path.name)
    for page_idx, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        metadata = {
            "source_file": file_path.name,
            "session": session_num,
            "page": page_idx
        }
        docs.append((page_text, metadata))
    return docs

def extract_text_from_pptx_per_slide(file_path):
    """
    Returns a list of (slide_text, metadata) for each slide in the PPTX.
    """
    pres = Presentation(file_path)
    docs = []
    # Optionally parse session number from filename
    session_num = parse_session_from_filename(file_path.name)
    for slide_idx, slide in enumerate(pres.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + " "
        
        metadata = {
            "source_file": file_path.name,
            "session": session_num,
            "slide": slide_idx
        }
        docs.append((slide_text.strip(), metadata))
    return docs

def parse_session_from_filename(filename):
    """
    Optional helper to extract a session number from file names like 'session_7.pdf'.
    Adjust as needed.
    """
    # Simple example: look for something like 'session_7'
    # Return int(7) if found, else None
    # This is just an example; adapt it to your real file naming convention
    import re
    match = re.search(r'session_(\d+)', filename.lower())
    return int(match.group(1)) if match else None

def process_files(input_folder):
    """
    Returns a list of (text, metadata) across all PPTX/PDF files in folder.
    """
    folder = Path(input_folder)
    all_docs = []  # will hold a list of (text, metadata) pairs
    
    for file_path in folder.glob("*"):
        if file_path.suffix.lower() == ".pdf":
            docs = extract_text_from_pdf_per_page(file_path)
            all_docs.extend(docs)
        elif file_path.suffix.lower() == ".pptx":
            docs = extract_text_from_pptx_per_slide(file_path)
            all_docs.extend(docs)
        else:
            print(f"Unsupported file format: {file_path}")
    
    return all_docs

def clean_text(content: str) -> str:
    """
    Same as before – cleans a single string.
    You might adapt it if you'd like to handle smaller strings per slide/page.
    """
    lines = content.splitlines()
    cleaned_lines = [
        line.strip()
        for line in lines
        if not line.isdigit() and len(line.strip()) > 10
    ]
    return "\n".join(cleaned_lines)
